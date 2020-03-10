# coding: UTF-8

from pptx import Presentation
from pdf2image import convert_from_path, convert_from_bytes
from pdf2image.exceptions import (
    PDFInfoNotInstalledError,
    PDFPageCountError,
    PDFSyntaxError
)
import os
import datetime
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

dt_name = datetime.datetime.now()
date = '%d-%d-%d' % (dt_name.year,dt_name.month,dt_name.day)

for filename in os.listdir('source_files/'):
    if os.path.splitext(filename)[1] == '.pdf':
        print("Creating %s" % filename)
        prs = Presentation()

        pages = convert_from_path('source_files/' + filename, 500)
        for index, page in enumerate(pages):
            #Save as 'jpg' in jpgs dir
            jpg_file = "jpgs/%s-(%d).jpg" % (filename,index)
            page.save(jpg_file, 'JPEG')

            #Get width/height of image
            image = Image.open(jpg_file)
            height = image.height
            width = image.width
            #Rotate 270 degrees if horizontal
            if height > width:
                adjusted = image.rotate(270, expand=True)
                adjusted.save(jpg_file)

            #Setup slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            left = top = 0
            slide.shapes.add_picture(jpg_file, left-0.1*prs.slide_width,top,height = prs.slide_height)


            #Create Top left box
            top_left_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.17), Inches(0.17), Inches(3.39), Inches(0.29)
            )
            #Fill Top left box
            fill = top_left_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(74,125,186)

            #Text Top left box
            text_frame = top_left_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = 'TC_%s_%s' % (os.path.splitext(filename)[0],date)

            font = run.font
            font.name = 'Arial'
            font.size = Pt(10)
            font.bold = True


            #Create Top right box
            tr_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(8.86), Inches(0.20), Inches(1.0), Inches(0.29)
            )
            tr_fill = tr_shape.fill
            tr_fill.solid()
            tr_fill.fore_color.rgb = RGBColor(74,125,186)
            tr_text_frame = tr_shape.text_frame
            tr_text_frame.clear()
            tr_p = tr_text_frame.paragraphs[0]
            tr_run = tr_p.add_run()
            tr_run.text = date
            tr_font = tr_run.font
            tr_font.name = 'Arial'
            tr_font.size = Pt(10)
            tr_font.bold = True

        prs.save('result/%s.pptx' % os.path.splitext(filename)[0])

    else:
        print("Skipping %s because it\'s not a pdf" % filename)

print("Saved to result directory")
